import copy
import re

from jinja2 import exceptions, Environment
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import _Cell

from template_pptx_jinja import pictures


class PPTXRendering:
    table_data_pattern = re.compile(r"\{\{\s*table\s*:\s*([^}\s]+)\s*\}\}")

    def __init__(self, input_path, data, output_path, env=None):
        self.input_path = input_path
        self.model = data['model']
        self.pictures = data.get('pictures')

        self.output_path = output_path
        if env is not None:
            self.env = env
        else:
            self.env = Environment()

        self.current_slide = None
        self.current_shape = None

        self.message_raw = []

    def process(self):
        ppt = Presentation(self.input_path)
        for slide in ppt.slides:
            self.current_slide = slide
            self._render_slide(slide)
        ppt.save(self.output_path)
        return "\n".join(self.message_raw)

    def _render_slide(self, slide):
        for shape in slide.shapes:
            self.current_shape = shape
            self._render_shape(shape)

    def _render_shape(self, shape):
        if shape.has_text_frame:
            self._render_text_frame(shape.text_frame)
        if shape.has_table:
            self._render_table(shape.table)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._render_picture(shape)

    def _render_picture(self, shape):
        for picture in self.pictures:
            if pictures.get_hash(picture) == shape.image.sha1:
                pictures.replace_img_slide(
                    self.current_slide, shape, self.pictures[picture]
                )

    def _render_text_frame(self, text_frame):
        for paragraph in text_frame.paragraphs:
            self._render_paragraph(paragraph)

    def _render_paragraph(self, paragraph):
        for run in paragraph.runs:
            self._render_run(run)

    def _render_table(self, table):
        self._prepare_table(table)
        for cell in table.iter_cells():
            self._render_cell(cell)

    def _prepare_table(self, table):
        for cell in table.iter_cells():
            matches = re.findall(self.table_data_pattern, cell.text)
            if matches:
                table_data_key = matches[0].strip()
                self._remove_row(table, 0)
                if table_data := self.model.get(table_data_key):
                    for i, row in enumerate(table_data):
                        self._add_row(table, i)
                    self._remove_row(table, 1)
            break

    def _render_cell(self, cell):
        self._render_text_frame(cell.text_frame)

    def _render_run(self, run):
        try:
            template = self.env.from_string(str(run.text))
            rendered = template.render(self.model)
        except exceptions.UndefinedError as error:
            error_text = f"{error.__class__.__name__}: {error}"
            self.message_raw.append(error_text)
        except exceptions.TemplateSyntaxError as error:
            error_text = (
                f"{error.__class__.__name__}: {error}\n"
                "you should re-write the whole {{}} tag"
            )
            self.message_raw.append(error_text)
        else:
            run.text = rendered

    @staticmethod
    def _add_row(table, row_index):
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])

        for cell_index, tc in enumerate(new_row.tc_lst):
            cell = _Cell(tc, new_row.tc_lst)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = f"{{{{ table_data[{row_index}][{cell_index}] }}}}"

        table._tbl.append(new_row)

    @staticmethod
    def _remove_row(table, row_index):
        row = list(table.rows)[row_index]
        table._tbl.remove(row._tr)

